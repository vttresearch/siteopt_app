from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT_PPTX = ROOT / "SiteOpt-Assistant-Achievement.pptx"

IMG_ASSISTANT = ROOT / "Screenshot from 2026-02-23 11-26-28.png"
IMG_PLOT = ROOT / "Screenshot from 2026-02-23 11-26-49.png"
IMG_COPILOT_VIS = ROOT / "Screenshot from 2026-02-23 13-19-46.png"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_title_slide(title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(24)


def add_image_slide(title: str, image_path: Path, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.7), Inches(1.2), height=Inches(4.8))

    textbox = slide.shapes.add_textbox(Inches(8.0), Inches(1.4), Inches(5.0), Inches(4.8))
    tf = textbox.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(22)


add_title_slide(
    "SiteOpt Assistant",
    "User-Focused Achievement Summary\nPurpose: Help users prepare inputs faster and understand data before running energy system studies.",
)

add_bullet_slide(
    "1. What Users Can Do Now",
    [
        "Ask for input changes in plain language.",
        "Ask questions about project data without manual searching.",
        "Request visualizations directly from chat.",
        "See generated plots in Data Editor → Assistant Plot.",
    ],
)

add_bullet_slide(
    "2. End-User Experience",
    [
        "Assistant is available directly in the web interface.",
        "Works in the currently active project context.",
        "Responses are readable and structured.",
        "Charts are shown where users already inspect data.",
    ],
)

add_bullet_slide(
    "3. Separation of Responsibilities",
    [
        "Web Assistant: input preparation and insight.",
        "Spine Toolbox: optimization/simulation execution.",
        "Clear workflow: prepare in web app, execute in Spine Toolbox, review results.",
    ],
)

add_image_slide(
    "4. Assistant in Daily Workflow",
    IMG_ASSISTANT,
    [
        "User asks naturally.",
        "Assistant responds with actionable guidance and summaries.",
    ],
)

add_image_slide(
    "5. Visualization in Data Editor",
    IMG_PLOT,
    [
        "Requested chart appears in Assistant Plot tab.",
        "Users stay in one workspace while preparing inputs.",
    ],
)

add_image_slide(
    "6. Copilot-Requested Data Visualization",
    IMG_COPILOT_VIS,
    [
        "User asks for a specific data view in natural language.",
        "Copilot generates a matching visualization for decision support.",
    ],
)

add_bullet_slide(
    "7. Practical Value for Users",
    [
        "Less time spent on repetitive manual edits.",
        "Faster confidence checks with quick plots and summaries.",
        "Lower risk of confusion from switching between many tools too early.",
        "Better collaboration: requests and outputs are easy to discuss.",
    ],
)

add_bullet_slide(
    "8. Message to Stakeholders",
    [
        "Usability improved without changing core simulation responsibility.",
        "Supports current workflow rather than replacing it.",
        "Creates a clear bridge between input preparation and model runs.",
    ],
)

add_bullet_slide("Thank You", ["Questions / Feedback"])

prs.save(str(OUT_PPTX))
print(f"Created: {OUT_PPTX}")
